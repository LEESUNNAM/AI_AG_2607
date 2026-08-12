"""
General-purpose .pptx slide deck builder using python-pptx.

Styling defaults match this project's other document scripts
(scripts/create_docx.py, scripts/create_xlsx.py): Malgun Gothic font,
a customizable main color applied to titles and accents, 16:9 widescreen.

Every slide is built on a blank layout with manually placed text boxes /
pictures rather than relying on the default placeholder theme — this keeps
Korean text rendering consistent (PowerPoint's default theme font falls
back to a Latin-only face for East Asian runs unless set explicitly) and
keeps every deck visually consistent with this project's other outputs.

Usage as a library:
    from create_pptx import PptxBuilder

    deck = PptxBuilder(main_color="1F4E79")
    deck.add_title_slide("Report Title", subtitle="2026 Analysis")
    deck.add_section_slide("1. Overview")
    deck.add_content_slide("Key Points", ["Point A", "Point B", "Point C"])
    deck.add_image_slide("Screenshot", "output2/playwright/example.png", caption="Source: ...")
    deck.add_table_slide("Data", headers=["Item", "Value"], rows=[["Alpha", "1"], ["Beta", "2"]])
    deck.save("output/deck.pptx")

Usage from the command line (creates a demo deck):
    python create_pptx.py output/demo.pptx --color 2E86AB --title "Sample Deck"
"""

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

DEFAULT_FONT = "맑은 고딕"
DEFAULT_MAIN_COLOR = "1F4E79"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.6)


def _parse_color(color: str | RGBColor) -> RGBColor:
    if isinstance(color, RGBColor):
        return color
    return RGBColor.from_string(color.lstrip("#").upper())


def _set_font(run, font_name: str, size_pt: float, color: RGBColor | None = None, bold: bool | None = None):
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    # python-pptx's font.name only sets the Latin typeface; Korean text needs
    # the East Asian typeface set explicitly or PowerPoint falls back to its default font.
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea is None:
        from pptx.oxml.ns import qn
        ea = rpr.makeelement(qn("a:ea"), {})
        rpr.append(ea)
    ea.set("typeface", font_name)


class PptxBuilder:
    def __init__(
        self,
        main_color: str = DEFAULT_MAIN_COLOR,
        font_name: str = DEFAULT_FONT,
        accent_color: str = "FFFFFF",
    ):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self.main_color = _parse_color(main_color)
        self.accent_color = _parse_color(accent_color)
        self.font_name = font_name
        self._blank_layout = self.prs.slide_layouts[6]

    def _new_slide(self):
        return self.prs.slides.add_slide(self._blank_layout)

    def _add_textbox(self, slide, left, top, width, height, text, size, color=None, bold=False, align=PP_ALIGN.LEFT, anchor=None):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        if anchor is not None:
            tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        _set_font(run, self.font_name, size, color=_parse_color(color) if color else None, bold=bold)
        return box

    def _fill_background(self, slide, color: RGBColor):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def add_title_slide(self, title: str, subtitle: str | None = None, footer: str | None = None):
        slide = self._new_slide()
        self._fill_background(slide, self.main_color)
        self._add_textbox(
            slide, MARGIN, Inches(2.7), SLIDE_WIDTH - 2 * MARGIN, Inches(1.5),
            title, size=40, color=self.accent_color, bold=True, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        if subtitle:
            self._add_textbox(
                slide, MARGIN, Inches(4.2), SLIDE_WIDTH - 2 * MARGIN, Inches(0.8),
                subtitle, size=20, color=self.accent_color, align=PP_ALIGN.CENTER,
            )
        if footer:
            self._add_textbox(
                slide, MARGIN, SLIDE_HEIGHT - Inches(0.9), SLIDE_WIDTH - 2 * MARGIN, Inches(0.5),
                footer, size=12, color=self.accent_color, align=PP_ALIGN.CENTER,
            )
        return self

    def add_section_slide(self, title: str):
        """Full-bleed divider slide for a new section, matching the title slide theme."""
        slide = self._new_slide()
        self._fill_background(slide, self.main_color)
        self._add_textbox(
            slide, MARGIN, Inches(3.2), SLIDE_WIDTH - 2 * MARGIN, Inches(1.2),
            title, size=32, color=self.accent_color, bold=True, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        return self

    def _add_slide_header(self, slide, title: str):
        self._add_textbox(
            slide, MARGIN, Inches(0.35), SLIDE_WIDTH - 2 * MARGIN, Inches(0.8),
            title, size=26, color=self.main_color, bold=True,
        )
        line = slide.shapes.add_connector(1, MARGIN, Inches(1.15), SLIDE_WIDTH - MARGIN, Inches(1.15))
        line.line.color.rgb = self.main_color
        line.line.width = Pt(1.5)

    def add_content_slide(self, title: str, bullets: list[str], notes: str | None = None):
        slide = self._new_slide()
        self._add_slide_header(slide, title)
        box = slide.shapes.add_textbox(MARGIN, Inches(1.5), SLIDE_WIDTH - 2 * MARGIN, SLIDE_HEIGHT - Inches(2.0))
        tf = box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ""
            run = p.add_run()
            run.text = f"•  {bullet}"
            _set_font(run, self.font_name, 18)
            p.space_after = Pt(14)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return self

    def add_image_slide(self, title: str, image_path: str, caption: str | None = None):
        slide = self._new_slide()
        self._add_slide_header(slide, title)
        area_top = Inches(1.5)
        area_height = SLIDE_HEIGHT - Inches(2.2) - (Inches(0.5) if caption else 0)
        area_width = SLIDE_WIDTH - 2 * MARGIN
        pic = slide.shapes.add_picture(image_path, 0, 0)
        scale = min(area_width / pic.width, area_height / pic.height, 1.0)
        pic.width = Emu(int(pic.width * scale))
        pic.height = Emu(int(pic.height * scale))
        pic.left = Emu(int((SLIDE_WIDTH - pic.width) / 2))
        pic.top = Emu(int(area_top + (area_height - pic.height) / 2))
        if caption:
            self._add_textbox(
                slide, MARGIN, SLIDE_HEIGHT - Inches(0.7), SLIDE_WIDTH - 2 * MARGIN, Inches(0.5),
                caption, size=12, color="595959", align=PP_ALIGN.CENTER,
            )
        return self

    def add_table_slide(self, title: str, headers: list[str], rows: list[list[str]]):
        slide = self._new_slide()
        self._add_slide_header(slide, title)
        n_rows, n_cols = len(rows) + 1, len(headers)
        table_shape = slide.shapes.add_table(
            n_rows, n_cols, MARGIN, Inches(1.6), SLIDE_WIDTH - 2 * MARGIN, Inches(0.5) * n_rows
        )
        table = table_shape.table
        for c, header in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.main_color
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = header
            _set_font(run, self.font_name, 14, color=self.accent_color, bold=True)
        for r, row in enumerate(rows, start=1):
            for c, value in enumerate(row):
                cell = table.cell(r, c)
                cell.text = ""
                p = cell.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = str(value)
                _set_font(run, self.font_name, 13)
        return self

    def add_quote_slide(self, quote: str, attribution: str | None = None):
        """Full-slide pull-quote, useful for highlighting a key line from a scraped article."""
        slide = self._new_slide()
        self._add_textbox(
            slide, Inches(1.2), Inches(2.3), SLIDE_WIDTH - Inches(2.4), Inches(2.2),
            f"“{quote}”", size=28, color=self.main_color, bold=True, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        if attribution:
            self._add_textbox(
                slide, Inches(1.2), Inches(4.6), SLIDE_WIDTH - Inches(2.4), Inches(0.6),
                f"— {attribution}", size=16, color="595959", align=PP_ALIGN.CENTER,
            )
        return self

    def save(self, path: str):
        out_path = Path(path)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(out_path)
        return out_path


def _demo(path: str, title: str, main_color: str):
    deck = PptxBuilder(main_color=main_color)
    deck.add_title_slide(title, subtitle="Generated by create_pptx.py")
    deck.add_section_slide("1. Overview")
    deck.add_content_slide("Key Points", ["First point", "Second point", "Third point"])
    deck.add_table_slide("Data", headers=["Item", "Value"], rows=[["Alpha", "1"], ["Beta", "2"], ["Gamma", "3"]])
    deck.add_quote_slide("This is a sample pull quote.", attribution="Demo Source")
    saved_path = deck.save(path)
    print(f"Saved: {saved_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate a demo .pptx file.")
    parser.add_argument("output", nargs="?", default="output/demo.pptx")
    parser.add_argument("--title", default="Sample Deck")
    parser.add_argument("--color", default=DEFAULT_MAIN_COLOR, help="Main color as a hex string, e.g. 2E86AB")
    args = parser.parse_args()
    _demo(args.output, args.title, args.color)
